import os
import json
import re
import tempfile
import asyncio
import hashlib
import mimetypes
from datetime import datetime
from fastapi import APIRouter, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openai import AsyncOpenAI
import shutil
from .minio_store import (
    minio_enabled,
    minio_upload_file,
    minio_object_key,
    minio_presigned_url
)
from .mysql_store import save_artifact_by_md5, _sanitize_config, load_app_config
from .models import ConfigModel

router = APIRouter()

def _file_md5(path: str) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()

def _extract_text(file_path: str, ext: str) -> str:
    text = ""
    try:
        if ext == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        elif ext in ['.docx', '.doc']:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in ['.txt', '.md']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
    except Exception as e:
        print(f"Error extracting text: {e}")
        return f"Error extracting text: {str(e)}"
    return text

@router.post("/analyze_file")
async def analyze_file(
    file: UploadFile = File(None),
    config: str = Form(...),
    file_url: str | None = Form(None),
):
    """
    Upload a file (PDF, DOCX, PPTX, TXT, MD) and analyze it using LLM.
    Returns a stream of progress updates and finally the result.
    """
    async def event_generator():
        tmp_path = None
        try:
            # Parse config
            config_data = load_app_config()
            if not config_data:
                raise RuntimeError("数据库配置为空，请先在前端设置并保存")
            try:
                req_cfg = json.loads(config) if config else {}
            except Exception:
                req_cfg = {}
            if isinstance(req_cfg, dict):
                for k, v in req_cfg.items():
                    if v is None:
                        continue
                    if isinstance(v, str) and not v.strip():
                        continue
                    config_data[k] = v

            conf = ConfigModel(
                openai_api_key=config_data.get("openai_api_key", ""),
                openai_base_url=config_data.get("openai_base_url", ""),
                llm_model=config_data.get("llm_model", ""),
                ocr_engine=str(config_data.get("ocr_engine", "vl")),
                vl_model=str(config_data.get("vl_model", "Pro/Qwen/Qwen2-VL-7B-Instruct")),
                vl_base_url=str(config_data.get("vl_base_url", "https://api.siliconflow.cn/v1")),
                vl_api_key=str(config_data.get("vl_api_key", "")),
                model_size=str(config_data.get("model_size", "base")),
                device=str(config_data.get("device", "cpu")),
                compute_type=str(config_data.get("compute_type", "int8")),
                capture_offset=float(config_data.get("capture_offset", 0.0)),
            )

            if not file and not (file_url or "").strip():
                yield json.dumps({"status": "error", "message": "请上传文件或提供 file_url"}) + "\n"
                return

            filename = ""
            ext = ""
            remote_content_type = ""

            def _infer_ext_from(url: str, content_type: str) -> str:
                from urllib.parse import urlparse

                path = urlparse(url).path
                guessed_ext = os.path.splitext(path)[1].lower()
                if guessed_ext:
                    return guessed_ext

                ct = (content_type or "").lower()
                if "pdf" in ct:
                    return ".pdf"
                if "word" in ct:
                    return ".docx"
                if "presentation" in ct or "powerpoint" in ct:
                    return ".pptx"
                if "markdown" in ct:
                    return ".md"
                if "text" in ct:
                    return ".txt"
                return ".txt"

            def _infer_filename_from(url: str, fallback_ext: str) -> str:
                from urllib.parse import urlparse

                base = os.path.basename(urlparse(url).path or "").strip()
                if base:
                    if base.lower().endswith(fallback_ext):
                        return base
                    if "." not in base and fallback_ext:
                        return f"{base}{fallback_ext}"
                    return base
                return f"imported{fallback_ext or ''}"

            if file:
                filename = file.filename
                ext = os.path.splitext(filename)[1].lower()
                yield json.dumps({"status": "progress", "message": f"Processing file: {filename}..."}) + "\n"
            else:
                url = (file_url or "").strip()
                yield json.dumps({"status": "progress", "message": "正在从 URL 下载文件..."}) + "\n"
                import aiohttp

                async with aiohttp.ClientSession() as session:
                    async with session.get(url) as resp:
                        if resp.status != 200:
                            yield json.dumps({"status": "error", "message": f"下载失败: {resp.status}"}) + "\n"
                            return
                        remote_content_type = (resp.headers.get("Content-Type", "") or "").split(";", 1)[0].strip()
                        content_bytes = await resp.read()
                        ext = _infer_ext_from(url, remote_content_type)
                        if ext in (".txt", "") and content_bytes.lstrip().startswith(b"%PDF"):
                            ext = ".pdf"
                        filename = _infer_filename_from(url, ext)
                        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                            tmp.write(content_bytes)
                            tmp_path = tmp.name

            # Save file temporarily to calculate MD5
            if file:
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                    content = await file.read()
                    tmp.write(content)
                    tmp_path = tmp.name
            
            # Calculate MD5
            file_md5 = await asyncio.to_thread(_file_md5, tmp_path)

            source_ref_path = ""
            if minio_enabled():
                # Upload to MinIO
                object_key = minio_object_key("files", file_md5, filename)
                await asyncio.to_thread(minio_upload_file, tmp_path, object_key)
                source_ref_path = object_key
            else:
                # Save locally for persistent access (e.g. PDF preview)
                storage_dir = os.path.join(os.getcwd(), "storage", "files")
                os.makedirs(storage_dir, exist_ok=True)
                # Use md5 + ext as filename to avoid duplicates and encoding issues
                local_rel_path = f"files/{file_md5}{ext}"
                local_abs_path = os.path.join(os.getcwd(), "storage", local_rel_path)
                
                # Copy from temp file to storage
                shutil.copy2(tmp_path, local_abs_path)
                source_ref_path = local_rel_path
            
            # Use tmp_path for text extraction (we have it locally anyway)
            # Clean up temp file later


            yield json.dumps({"status": "progress", "message": "Extracting text content..."}) + "\n"
            
            # Extract text
            # If we uploaded to MinIO, we still have tmp_path. If we saved to fast_output, we have saved_path.
            # But we might have skipped saving to fast_output.
            # So let's rely on tmp_path for extraction which is always available here.
            text_content = await asyncio.to_thread(_extract_text, tmp_path, ext)
            
            if not text_content.strip():
                 yield json.dumps({"status": "error", "message": "Failed to extract text or file is empty."}) + "\n"
                 return

            yield json.dumps({"status": "progress", "message": "Analyzing content with LLM..."}) + "\n"

            # Call LLM
            client = AsyncOpenAI(api_key=conf.openai_api_key, base_url=conf.openai_base_url)
            
            system_prompt = (
                "You are a helpful assistant that summarizes documents. "
                "Please provide a comprehensive summary of the following content in Markdown format. "
                "Structure it with '## 摘要', '## 关键点', and '## 结论'. "
                "If the document is technical or educational, try to extract key concepts. "
                "IMPORTANT: You MUST answer in Simplified Chinese (简体中文), NO English allowed in the summary text. "
                "Structure it with '## 标题', '## 摘要', '## 关键点', and '## 结论'. "
                "At the very end, ON A NEW LINE, output a section '## Tags' followed by 3-5 keywords separated by commas. "
                "Example:\n## 标题\n标题示例\n\n## 摘要\n...\n\n## Tags\n关键词1, 关键词2, 关键词3"
            )
            
            # Truncate text if too long (simple heuristic, can be improved)
            max_chars = 100000 
            truncated_text = text_content[:max_chars]
            if len(text_content) > max_chars:
                truncated_text += "\n...(content truncated)..."

            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Here is the document content:\n\n{truncated_text}"}
            ]

            response = await client.chat.completions.create(
                model=conf.llm_model,
                messages=messages,
                stream=True
            )

            full_report = ""
            async for chunk in response:
                if chunk.choices and chunk.choices[0].delta.content:
                    content = chunk.choices[0].delta.content
                    full_report += content

            def _sanitize_title(s: str) -> str:
                t = (s or "").strip()
                t = re.sub(r"\s+", " ", t)
                t = t.strip(" \t\r\n-—:：#")
                return t[:80]

            ai_title = ""
            title_match = re.search(r"^##\s*标题\s*\n+(.+?)(?:\n|$)", full_report, re.MULTILINE)
            if title_match:
                ai_title = _sanitize_title(title_match.group(1))

            tags = []
            print(f"[DEBUG] Full report length: {len(full_report)}")
            tags_match = re.search(
                r"^##\s*(?:Tags|标签|关键标签|关键词)\s*\n+(.*?)(?=^\s*##\s+|\Z)",
                full_report,
                re.IGNORECASE | re.DOTALL | re.MULTILINE,
            )

            content = ""
            if tags_match:
                content = (tags_match.group(1) or "").strip()
            else:
                lines = full_report.strip().split("\n")
                tail = "\n".join(lines[-20:])
                tags_match2 = re.search(r"(?:Tags|标签|关键词)\s*(?::|：)\s*(.*)", tail, re.IGNORECASE)
                if tags_match2:
                    content = (tags_match2.group(1) or "").strip()

            if content:
                print(f"[DEBUG] Found tags content: {content}")
                content = re.sub(r"^\s*[-*•]\s*", "", content, flags=re.MULTILINE)
                raw_tags = re.split(r"[,，;；、\n]+", content)
                cleaned = []
                for t in raw_tags:
                    s = str(t or "").strip()
                    s = s.lstrip("#").strip()
                    if not s:
                        continue
                    if s.startswith("##"):
                        continue
                    cleaned.append(s)
                uniq = []
                seen = set()
                for t in cleaned:
                    if t in seen:
                        continue
                    seen.add(t)
                    uniq.append(t)
                tags = uniq[:5]
                print(f"[DEBUG] Extracted tags: {tags}")
            else:
                print("[DEBUG] No tags found in report")

            # Construct response data similar to video analysis structure for frontend compatibility
            mime_type = ""
            if file and file.content_type:
                mime_type = (file.content_type or "").strip()
            elif remote_content_type:
                mime_type = remote_content_type

            guessed, _ = mimetypes.guess_type(filename)
            if not mime_type or mime_type.lower() == "application/octet-stream":
                mime_type = (guessed or "").strip() or "application/octet-stream"
            if ext == ".md":
                mime_type = "text/markdown"
            elif ext == ".txt" and (not guessed):
                mime_type = "text/plain"

            title_preference = str(config_data.get("title_preference") or "").strip().lower()
            final_title = filename
            if title_preference != "filename" and ai_title:
                if ext and not ai_title.lower().endswith(ext):
                    final_title = f"{ai_title}{ext}"
                else:
                    final_title = ai_title

            data = {
                "title": final_title,
                "summary": full_report[:200] + "...",
                "segments": [], 
                "mind_map": {}, 
                "raw_transcript": text_content,
                "media_type": "document",
                "mime_type": mime_type,
                "tags": tags,
            }

            # Save to MySQL
            try:
                meta = {
                    "config": _sanitize_config(config_data),
                    "original_filename": filename
                }
                if not file and (file_url or "").strip():
                    meta["original_url"] = (file_url or "").strip()

                # Save asset and ai_analysis
                await asyncio.to_thread(save_artifact_by_md5,
                    video_md5=file_md5,
                    media_type="document",
                    asset_type="document",
                    mime_type=mime_type,
                    display_name=final_title,
                    source_kind="minio" if minio_enabled() else "local",
                    source_ref=source_ref_path,
                    meta=meta,
                    artifact_type="ai_analysis",
                    content_json=data
                )
                # Save report
                await asyncio.to_thread(save_artifact_by_md5,
                    video_md5=file_md5,
                    media_type="document",
                    asset_type="document",
                    meta=meta,
                    artifact_type="report_markdown",
                    content_text=full_report
                )
                # Save transcript
                await asyncio.to_thread(save_artifact_by_md5,
                    video_md5=file_md5,
                    media_type="document",
                    asset_type="document",
                    source_kind="minio" if minio_enabled() else "local",
                    source_ref=source_ref_path,
                    meta=meta,
                    artifact_type="raw_transcript",
                    content_text=text_content
                )
            except Exception as e:
                print(f"Error saving to MySQL: {e}")
            
            yield json.dumps({
                "status": "success", 
                "report": full_report,
                "data": data,
                "video_url": "", 
                "video_md5": file_md5,
                "message": "Analysis completed!"
            }) + "\n"

        except Exception as e:
            import traceback
            traceback.print_exc()
            yield json.dumps({"status": "error", "message": str(e)}) + "\n"
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except:
                    pass

    return StreamingResponse(event_generator(), media_type="text/event-stream")

@router.post("/import_file_url")
async def import_file_url(data: dict):
    url = data.get("url")
    if not url:
        return {"error": "URL is required"}
    
    # Simple download logic (supports PDF, DOCX, PPTX, TXT, MD)
    import aiohttp
    
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url) as resp:
                if resp.status != 200:
                     return {"error": f"Failed to download file: {resp.status}"}
                
                content = await resp.read()
                
                # Determine extension from URL or Content-Type
                from urllib.parse import urlparse
                path = urlparse(url).path
                ext = os.path.splitext(path)[1].lower()
                if not ext:
                    # Try to guess from content-type
                    ct = resp.headers.get('Content-Type', '').lower()
                    if 'pdf' in ct: ext = '.pdf'
                    elif 'word' in ct: ext = '.docx'
                    elif 'presentation' in ct or 'powerpoint' in ct: ext = '.pptx'
                    elif 'text' in ct: ext = '.txt'
                    elif 'markdown' in ct: ext = '.md'
                    else: ext = '.txt' # Default
                
                # Extract text
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                    tmp.write(content)
                    tmp_path = tmp.name
                
                text = await asyncio.to_thread(_extract_text, tmp_path, ext)
                
                # Clean up temp file
                try:
                    os.remove(tmp_path)
                except:
                    pass
                
                return {
                    "status": "success",
                    "file_url": url, # Return original URL for PDF iframe if applicable
                    "file_content": text,
                    "file_type": ext,
                    "message": "File imported successfully"
                }
    except Exception as e:
        return {"error": str(e)}
