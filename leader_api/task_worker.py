import asyncio
import logging
import traceback
import json
import os
import re
import tempfile
import uuid
from urllib.parse import urlparse, unquote
import shutil
from .task_store import tasks, queue, TaskStatus
from .mysql_store import load_app_config
from .minio_store import (
    extract_job_id_from_object_key,
    minio_download_to_file,
    minio_enabled,
    minio_object_exists,
    minio_object_etag_md5,
    minio_object_key,
    minio_presigned_url,
    minio_upload_tree,
    normalize_video_object_key,
    minio_remove_folder,
)
from .models import AnalyzePathRequest, ConfigModel
from fast_video_summary import FastVideoAnalyzer
from .mysql_store import save_artifact_by_md5, _sanitize_config
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openai import AsyncOpenAI

# Re-use the md5 and download logic from analyze.py
# Ideally this should be in a shared utility module
import hashlib
import urllib.request
try:
    import yt_dlp
except ImportError:
    yt_dlp = None

def _file_md5(path: str) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()

def _download_url_to_file(url: str, path: str) -> None:
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=60) as resp, open(path, "wb") as out:
        while True:
            chunk = resp.read(1024 * 1024)
            if not chunk:
                break
            out.write(chunk)

def _download_video_smart(url: str, output_path: str, progress_callback=None) -> str:
    """
    Download video using yt-dlp if available and appropriate, otherwise direct download.
    Returns the actual path of the downloaded file.
    """
    use_ytdlp = False
    if yt_dlp:
        # Check for common video sites
        if any(x in url for x in ["bilibili.com", "youtube.com", "youtu.be", "vimeo.com"]):
            use_ytdlp = True
            
    if use_ytdlp:
        # Use yt-dlp
        base_path = os.path.splitext(output_path)[0]
        
        ydl_opts = {
            'format': 'bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best',
            'outtmpl': base_path + '.%(ext)s',
            'quiet': True,
            'no_warnings': True,
        }
        
        if progress_callback:
            def hook(d):
                if d['status'] == 'downloading':
                    try:
                        p = d.get('_percent_str', '0%').replace('%','')
                        progress_callback(float(p))
                    except:
                        pass
            ydl_opts['progress_hooks'] = [hook]

        try:
            print(f"Attempting yt-dlp download for: {url}")
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)
                if 'entries' in info:
                    info = info['entries'][0]
                
                ext = info.get('ext', 'mp4')
                final_path = f"{base_path}.{ext}"
                
                # Check likely candidates if exact match missing
                if not os.path.exists(final_path):
                    import glob
                    candidates = glob.glob(f"{base_path}.*")
                    if candidates:
                        final_path = candidates[0]
                
                if os.path.exists(final_path):
                    return final_path
        except Exception as e:
            print(f"yt-dlp download failed: {e}, falling back to direct download")

    # Direct download
    _download_url_to_file(url, output_path)
    return output_path

def _extract_text(file_path: str, ext: str) -> str:
    text = ""
    try:
        if ext == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        elif ext in ['.docx', '.doc']:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in ['.txt', '.md']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
    except Exception as e:
        print(f"Error extracting text: {e}")
        return f"Error extracting text: {str(e)}"
    return text

def _load_db_config() -> dict:
    cfg = load_app_config()
    if not isinstance(cfg, dict) or not cfg:
        raise RuntimeError("数据库配置为空，请先在前端设置并保存")
    return cfg

async def worker_loop():
    print(">>> Background worker started and waiting for tasks...")
    while True:
        try:
            # print("DEBUG: Worker waiting for queue...")
            task_id = await queue.get()
            print(f">>> Worker received task: {task_id}")
            task = tasks.get(task_id)
            
            if not task:
                print(f"DEBUG: Task {task_id} not found in store")
                queue.task_done()
                continue
                
            if task.status == TaskStatus.CANCELLED:
                print(f"DEBUG: Task {task_id} was cancelled")
                queue.task_done()
                continue
                
            if task.status != TaskStatus.PENDING:
                print(f"DEBUG: Task {task_id} status is {task.status}, skipping")
                queue.task_done()
                continue

            print(f">>> Starting processing task {task_id} for URL: {task.url}")
            task.status = TaskStatus.RUNNING
            task.progress = 5
            task.logs.append("Task started.")
            task.save()
            
            try:
                if task.type == "file" or task.type == "document":
                    await process_file_task(task)
                else:
                    await process_video_task(task)
                
                if task.status != TaskStatus.CANCELLED:
                    task.status = TaskStatus.COMPLETED
                    task.progress = 100
                    task.logs.append("Task completed successfully.")
                    task.save()
            except Exception as e:
                print(f"Task {task_id} failed: {e}")
                traceback.print_exc()
                task.status = TaskStatus.FAILED
                task.error = str(e)
                task.logs.append(f"Error: {str(e)}")
                task.save()
            finally:
                queue.task_done()
                
        except asyncio.CancelledError:
            print("Worker cancelled.")
            break
        except Exception as e:
            print(f"Worker error: {e}")
            await asyncio.sleep(1)

async def process_file_task(task):
    task.logs.append("Downloading file...")
    url = task.url
    task_cfg = _load_db_config()
    if isinstance(task.config, dict) and task.config:
        for k, v in task.config.items():
            if v is None:
                continue
            if isinstance(v, str) and not v.strip() and k != "title_preference":
                continue
            task_cfg[k] = v
    title_preference = str(task_cfg.get("title_preference") or "").strip().lower()
    
    # Create temp file
    parsed = urlparse(url)
    url_path = parsed.path if parsed.scheme in ("http", "https") else url
    ext = os.path.splitext(url_path)[1].lower() if "." in url_path else ""
    if ext:
        ext = ext.split("?", 1)[0].split("#", 1)[0]
    if not ext:
        ext = ".txt"
        
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp_path = tmp.name
    
    try:
        # Download
        if minio_enabled() and not url.startswith("http"):
             # Assume object key
             minio_download_to_file(url, tmp_path)
        else:
             _download_url_to_file(url, tmp_path)
             
        task.progress = 20
        task.logs.append("File downloaded. Calculating MD5...")
        
        # MD5
        file_md5 = _file_md5(tmp_path)
        
        # Extract Text
        task.progress = 30
        task.logs.append("Extracting text...")
        task.save()
        text_content = await asyncio.to_thread(_extract_text, tmp_path, ext)
        
        if not text_content.strip():
            raise Exception("Failed to extract text or file is empty.")
            
        task.progress = 50
        task.logs.append("Analyzing content with LLM...")
        task.save()
        
        # LLM Analysis
        co = task_cfg.get("capture_offset", None)
        config = ConfigModel(
            openai_api_key=(task_cfg.get("openai_api_key") or ""),
            openai_base_url=(task_cfg.get("openai_base_url") or ""),
            llm_model=(task_cfg.get("llm_model") or ""),
            ocr_engine=(task_cfg.get("ocr_engine") or "vl"),
            vl_model=(task_cfg.get("vl_model") or "Pro/Qwen/Qwen2-VL-7B-Instruct"),
            vl_base_url=(task_cfg.get("vl_base_url") or "https://api.siliconflow.cn/v1"),
            vl_api_key=(task_cfg.get("vl_api_key") or ""),
            model_size=(task_cfg.get("model_size") or "base"),
            device=(task_cfg.get("device") or "cpu"),
            compute_type=(task_cfg.get("compute_type") or "int8"),
            capture_offset=float(co) if co is not None and str(co).strip() != "" else 0.0,
        )
        
        client = AsyncOpenAI(api_key=config.openai_api_key, base_url=config.openai_base_url)
        system_prompt = (
            "你是一个擅长总结文档的助手。请用 Markdown 对以下内容做结构化总结，必须使用简体中文。"
            "结构必须包含：'## 标题'、'## 摘要'、'## 关键点'、'## 结论'。"
            "'## 标题' 下面只允许一行短标题。"
            "在文章末尾另起一段输出 '## 标签'，给出 3-5 个关键词，用逗号分隔。"
        )
        
        max_chars = 100000 
        truncated_text = text_content[:max_chars]
        if len(text_content) > max_chars:
            truncated_text += "\n...(content truncated)..."

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Here is the document content:\n\n{truncated_text}"}
        ]

        response = await client.chat.completions.create(
            model=config.llm_model,
            messages=messages,
            stream=False # Simplify for background task
        )
        
        full_report = response.choices[0].message.content
        
        task.progress = 90
        task.logs.append("Saving results...")
        
        # Save artifacts
        filename = unquote(os.path.basename(url_path.replace("\\", "/"))).strip()
        if not filename:
            filename = f"document{ext}"
        def _sanitize_title(s: str) -> str:
            t = (s or "").strip()
            t = re.sub(r"\s+", " ", t)
            t = t.strip(" \t\r\n-—:：#")
            return t[:80]

        display_name = filename
        if title_preference == "ai":
            m = re.search(r"##\s*标题\s*(.*?)(?=\n##|\Z)", full_report or "", re.DOTALL | re.IGNORECASE)
            ai_title = _sanitize_title(m.group(1)) if m else ""
            if ai_title:
                display_name = ai_title

        def _extract_tags(text: str) -> list[str]:
            if not text:
                return []
            m = re.search(
                r"^##\s*(?:Tags|标签|关键标签|关键词)\s*\n+(.*?)(?=^\s*##\s+|\Z)",
                text,
                re.IGNORECASE | re.DOTALL | re.MULTILINE,
            )
            content = ""
            if m:
                content = (m.group(1) or "").strip()
            else:
                lines = (text or "").strip().split("\n")
                tail = "\n".join(lines[-20:])
                m2 = re.search(r"(?:Tags|标签|关键词)\s*(?::|：)\s*(.*)", tail, re.IGNORECASE)
                if m2:
                    content = (m2.group(1) or "").strip()
            if not content:
                return []
            content = re.sub(r"^\s*[-*•]\s*", "", content, flags=re.MULTILINE)
            parts = re.split(r"[,，;；、\n]+", content)
            tags = []
            for p in parts:
                t = str(p or "").strip()
                t = t.lstrip("#").strip()
                if not t:
                    continue
                if t.startswith("##"):
                    continue
                tags.append(t)
            uniq = []
            seen = set()
            for t in tags:
                if t in seen:
                    continue
                seen.add(t)
                uniq.append(t)
            return uniq[:5]

        tags = _extract_tags(full_report)

        meta = {
            "config": config.model_dump(),
            "original_filename": filename,
            "job_id": task.id
        }
        
        # Detect MIME type
        import mimetypes
        mime_type, _ = mimetypes.guess_type(url_path)
        if ext == ".md":
            mime_type = "text/markdown"
        if not mime_type:
            mime_type = "application/octet-stream"

        await asyncio.to_thread(save_artifact_by_md5,
            video_md5=file_md5,
            media_type="document",
            asset_type="document",
            mime_type=mime_type, 
            display_name=display_name,
            source_kind="minio" if minio_enabled() and not url.startswith("http") else "url",
            source_ref=url,
            meta=meta,
            artifact_type="ai_analysis",
            artifact_version=1,
            content_json={
                "title": display_name,
                "summary": full_report[:200] + "...",
                "tags": tags,
                "raw_transcript": text_content,
                "media_type": "document"
            }
        )
        
        # Save Report
        await asyncio.to_thread(save_artifact_by_md5,
            video_md5=file_md5,
            media_type="document",
            asset_type="document",
            source_kind="minio" if minio_enabled() and not url.startswith("http") else "url",
            source_ref=url,
            meta=meta,
            artifact_type="report_markdown",
            artifact_version=1,
            content_text=full_report
        )

        await asyncio.to_thread(save_artifact_by_md5,
            video_md5=file_md5,
            media_type="document",
            asset_type="document",
            source_kind="minio" if minio_enabled() and not url.startswith("http") else "url",
            source_ref=url,
            meta=meta,
            artifact_type="raw_transcript",
            artifact_version=1,
            content_text=text_content
        )
        
        task.result = {
            "md5": file_md5,
        }

    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)

async def process_video_task(task):
    # This is a simplified version of analyze_video_by_path adapted for background execution
    video_ref = task.url
    
    # Create default config
    task_cfg = _load_db_config()
    if isinstance(task.config, dict) and task.config:
        for k, v in task.config.items():
            if v is None:
                continue
            if isinstance(v, str) and not v.strip() and k != "title_preference":
                continue
            task_cfg[k] = v
    task_cfg_log = None
    try:
        task_cfg_log = _sanitize_config(task_cfg)
    except Exception:
        task_cfg_log = {}
    print(f">>> [DEBUG] Worker received task config: {task_cfg_log}")
    
    def get_cfg(key, default):
        val = task_cfg.get(key)
        if val is not None and str(val).strip():
            print(f">>> [DEBUG] Using task_cfg for {key}")
            return val
        print(f">>> [DEBUG] Using default for {key}: {default}")
        return default

    co = task_cfg.get("capture_offset", None)
    title_preference = str(task_cfg.get("title_preference") or "").strip().lower()
    
    config = ConfigModel(
        openai_api_key=get_cfg("openai_api_key", ""),
        openai_base_url=get_cfg("openai_base_url", "https://api.openai.com/v1"),
        llm_model=get_cfg("llm_model", "gpt-3.5-turbo"),
        ocr_engine=get_cfg("ocr_engine", "vl"),
        vl_model=get_cfg("vl_model", "Pro/Qwen/Qwen2-VL-7B-Instruct"),
        vl_base_url=get_cfg("vl_base_url", "https://api.siliconflow.cn/v1"),
        vl_api_key=get_cfg("vl_api_key", ""),
        model_size=get_cfg("model_size", "base"),
        device=get_cfg("device", "cpu"),
        compute_type=get_cfg("compute_type", "int8"),
        capture_offset=float(co) if co is not None and str(co).strip() != "" else 0.5,
    )
    
    print(f">>> [DEBUG] Final ConfigModel: {config.model_dump()}")
    
    task.logs.append("Initializing analysis...")
    
    # DEBUG: Print API Key source (masked)
    api_key = config.openai_api_key
    masked_key = f"{api_key[:8]}...{api_key[-4:]}" if api_key and len(api_key) > 12 else "EMPTY/INVALID"
    print(f"DEBUG: Using API Key: {masked_key}")
    print(f"DEBUG: Base URL: {config.openai_base_url}")
    print(f"DEBUG: Model: {config.llm_model}")
    task.logs.append(f"Config: Model={config.llm_model}, Key={masked_key}")

    parsed = urlparse(video_ref)
    job_id = ""
    video_source = ""
    video_md5 = ""
    
    # Setup job context
    if parsed.scheme in ("http", "https"):
        job_id = str(uuid.uuid4())
    else:
        # MinIO path support (if needed)
        job_id = str(uuid.uuid4())

    # Use persistent directory instead of temporary one to ensure artifacts (video, images) survive
    # fast_output is mounted as /static in app.py (REMOVED by user request)
    # If MinIO is enabled, we can use a temporary directory because artifacts will be uploaded to MinIO.
    use_minio = minio_enabled()
    
    # Always use temporary directory for processing
    # If MinIO is not enabled, artifacts will be lost after processing (as fast_output is removed)
    temp_dir_obj = tempfile.TemporaryDirectory(prefix="leader-task-")
    base_output_dir = temp_dir_obj.name
    # We don't need "tasks" subdir really, but to keep structure similar:
    output_dir = os.path.join(base_output_dir, job_id)
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        # Download
        print(f">>> Task {task.id}: Starting download for {video_ref}")
        task.logs.append("Downloading video...")
        local_source_target = os.path.join(output_dir, "source.mp4")
        
        def download_progress(p):
            # p is 0-100 float
            # Map 0-100 download to task progress 5-15
            try:
                task.progress = 5 + int(float(p) * 0.1)
            except:
                pass
        
        try:
            local_source = await asyncio.to_thread(_download_video_smart, video_ref, local_source_target, download_progress)
            print(f">>> Task {task.id}: Download finished: {local_source}")
        except Exception as e:
            print(f"!!! Task {task.id} download failed: {e}")
            raise RuntimeError(f"Failed to download video: {e}")
            
        task.progress = 15
        task.save()
        
        # MD5
        print(f">>> Task {task.id}: Calculating MD5...")
        task.logs.append("Calculating MD5...")
        if not os.path.exists(local_source):
             raise RuntimeError(f"Video file not found at {local_source}")
             
        video_md5 = await asyncio.to_thread(_file_md5, local_source)
        print(f">>> Task {task.id}: MD5 = {video_md5}")
        video_source = local_source
        task.progress = 20
        task.save()
        
        # Analyze
        print(f">>> Task {task.id}: Starting analyzer...")
        task.logs.append("Starting analysis model...")
        
        def progress_cb(msg):
            print(f">>> Task {task.id} Progress: {msg}")
            task.logs.append(msg)
            # Rough mapping of messages to progress
            if "Transcribing" in msg or "转录" in msg:
                task.progress = 30
            elif "LLM" in msg:
                task.progress = 60
            elif "Summary" in msg or "摘要" in msg:
                task.progress = 80
            elif "Report" in msg or "报告" in msg:
                task.progress = 90
            # Don't save on every callback to avoid DB spam
        
        analyzer = FastVideoAnalyzer(
            video_source,
            output_dir,
            config=config.model_dump(),
            progress_callback=progress_cb,
        )
        
        ai_analysis = await asyncio.to_thread(analyzer.run)
        
        if not ai_analysis:
            raise RuntimeError("Analysis returned no results")
            
        task.logs.append("Uploading results...")
        
        # Upload to MinIO (Optional)
        if use_minio:
            minio_upload_tree(output_dir, minio_object_key("outputs", job_id))
        
        # Read report
        report_path = os.path.join(output_dir, "final_report.md")
        report_content = ""
        if os.path.exists(report_path):
            with open(report_path, "r", encoding="utf-8") as f:
                report_content = f.read()

        # Fix image links
        # If MinIO enabled, use MinIO URLs (relative/object_key). Otherwise use local static URLs.
        
        for seg in ai_analysis.get("segments", []):
            if "image_path" in seg:
                # image_path from analyzer is relative to output_dir, e.g., "images/keyframe_00_12s.jpg"
                # We need to ensure it is correct.
                
                # Normalize path separators
                rel_path = seg["image_path"].replace("\\", "/")
                
                if use_minio:
                    image_object_key = minio_object_key("outputs", job_id, rel_path)
                    seg["image_object_key"] = image_object_key
                    # Do NOT bake in presigned URL to report or analysis to avoid expiration
                    # Leave image_url as relative path or empty, analyze.py will generate it
                    seg["image_url"] = rel_path 
                    # Do NOT replace in report_content if using MinIO, analyze.py will handle relative paths
                else:
                    # No local storage (fast_output removed)
                    seg["image_url"] = ""
        
        # Save to MySQL
        task.logs.append("Saving to database...")
        original_filename = unquote(os.path.basename(parsed.path) or "video.mp4")
        def _sanitize_title(s: str) -> str:
            t = (s or "").strip()
            t = re.sub(r"\s+", " ", t)
            t = t.strip(" \t\r\n-—:：#")
            return t[:80]

        display_name = original_filename
        ai_title = ""
        try:
            ai_title = _sanitize_title(str(ai_analysis.get("title") or ""))
        except Exception:
            ai_title = ""
        if not ai_title:
            m = re.search(r"##\s*标题\s*(.*?)(?=\n##|\Z)", report_content or "", re.DOTALL | re.IGNORECASE)
            ai_title = _sanitize_title(m.group(1)) if m else ""
        if title_preference == "ai" and ai_title:
            display_name = ai_title
        try:
            ai_analysis["title"] = display_name
            if ai_title:
                ai_analysis["ai_title"] = ai_title
        except Exception:
            pass

        meta = {
            "config": config.model_dump(),
            "original_filename": original_filename,
            "ai_title": ai_title,
            "job_id": job_id
        }
        
        # Determine video source ref for database
        if use_minio:
            saved_source_kind = "minio"
            saved_source_ref = minio_object_key("outputs", job_id, "source.mp4")
        else:
            # Without fast_output/static, we don't have a valid local URL
            saved_source_kind = "none"
            saved_source_ref = ""
        
        # If original was URL, maybe keep it? But we want to play the downloaded file.
        # The frontend Video.vue or AiVideoSummary.vue expects to play something.
        
        await asyncio.to_thread(save_artifact_by_md5,
            video_md5=video_md5,
            media_type="video",
            asset_type="video",
            mime_type="video/mp4",
            display_name=display_name,
            source_kind=saved_source_kind,
            source_ref=saved_source_ref,
            meta=meta,
            artifact_type="ai_analysis",
            content_json=ai_analysis
        )
        
        # Save report
        await asyncio.to_thread(save_artifact_by_md5,
            video_md5=video_md5,
            # Do not overwrite asset-level properties with report details
            # media_type="text",
            # asset_type="report_markdown",
            # mime_type="text/markdown", 
            # display_name="Final Report",
            # source_kind="generated",
            # source_ref=f"tasks/{job_id}/final_report.md",
            meta=meta, # Keep/Update asset meta with job_id
            artifact_type="report_markdown",
            content_text=report_content,
            artifact_meta=meta # Also store meta in artifact
        )

        # Extract and save captured frames
        captured_frames = []
        for seg in ai_analysis.get("segments", []):
            if "image_path" in seg:
                captured_frames.append({
                    "url": seg.get("image_url", ""),
                    "timestamp": seg.get("timestamp", 0),
                    "object_key": seg.get("image_object_key", "")
                })
        
        if captured_frames:
            await asyncio.to_thread(save_artifact_by_md5,
                video_md5=video_md5,
                media_type="video",
                asset_type="video",
                source_kind=saved_source_kind,
                source_ref=saved_source_ref,
                meta=meta,
                artifact_type="captured_frames",
                content_json=captured_frames
            )
        
        task.result = {
            "md5": video_md5,
            "job_id": job_id
        }
    except Exception as e:
        task.logs.append(f"Processing error: {str(e)}")
        raise e
    finally:
        # Cleanup if using temp dir
        if temp_dir_obj:
            try:
                temp_dir_obj.cleanup()
            except:
                pass
